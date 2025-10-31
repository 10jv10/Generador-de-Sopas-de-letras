import streamlit as st # Streamlit se mantiene aquí para los st.error/warning dentro de las funciones de lógica/dibujo
import pandas as pd
import random
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from io import BytesIO
from reportlab.pdfgen import canvas
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.lib.colors import black
from reportlab.lib.units import inch
from reportlab.lib.pagesizes import letter # Usaremos 'letter' como base

# --- CONSTANTES GLOBALES DE DISEÑO ---

# Diccionario de tamaños KDP (ancho, alto en cm)
TAMAÑOS_KDP = {
    '5" x 8" (12.7 x 20.32 cm)': (12.7, 20.32),
    '5.06" x 7.81" (12.85 x 19.84 cm)': (12.85, 19.84),
    '5.25" x 8" (13.34 x 20.32 cm)': (13.34, 20.32),
    '5.5" x 8.5" (13.97 x 21.59 cm)': (13.97, 21.59),
    '6" x 9" (15.24 x 22.86 cm)': (15.24, 22.86),
    '6.14" x 9.21" (15.6 x 23.39 cm)': (15.6, 23.39),
    '6.69" x 9.61" (16.99 x 24.41 cm)': (16.99, 24.41),
    '7" x 10" (17.78 x 25.4 cm)': (17.78, 25.4),
    '8" x 10" (20.32 x 25.4 cm)': (20.32, 25.4),
    '8.25" x 6" (20.96 x 15.24 cm)': (20.96, 15.24),
    '8.25" x 8.25" (20.96 x 20.96 cm)': (20.96, 20.96),
    '8.27" x 11.69" (21.01 x 29.69 cm)': (21.01, 29.69),
    '8.5" x 8.5" (21.59 x 21.59 cm)': (21.59, 21.59),
    '8.5" x 11" (21.59 x 27.94 cm)': (21.59, 27.94)
}

# --- CONSTANTES DE FUENTES ---
FONT_NAME_REGULAR = "RobotoMono-Regular"
FONT_NAME_BOLD = "RobotoMono-Bold"
TTF_FILE_REGULAR = "RobotoMono-Regular.ttf"
TTF_FILE_BOLD = "RobotoMono-Bold.ttf"

# --- CONSTANTES DE DISEÑO Y LAYOUT ---
PADDING = 0.75 * inch            # Margen de seguridad interior (sangrado)
TITLE_FONT_SIZE = 18             # Tamaño fijo para el título de la sopa
WORDS_FONT_SIZE = 11             # Tamaño fijo para la lista de palabras
SPACE_AFTER_TITLE = 0.2 * inch   # Espacio entre título y lista de palabras
INTERLINEADO_PALABRAS = 1.3      # Multiplicador de altura de línea
PROPORCION_LETRA_CUADRICULA = 0.7  # Letra ocupa el 70% de la celda
NUM_COLUMNAS_PALABRAS = 4        # Columnas para la lista de palabras
SOLUCIONES_POR_PAGINA = 4        # Cuadrícula de 2x2 en páginas de solución
LETRAS_ALFABETO = 'ABCDEFGHIJKLMNOPQRSTUVWXYZ'

# Constantes para el recuadro de palabras
BOX_PADDING = 0.15 * inch        # Relleno interno del recuadro
BOX_CORNER_RADIUS = 0.1 * inch   # Radio de las esquinas redondeadas

# --- FUNCIONES DE LÓGICA ---

def cm_to_pt(cm):
    """
    Convierte una medida de centímetros (cm) a puntos (pt) de ReportLab.
    """
    return cm * 28.3465

def crear_sopa_letras(palabras, dimension=20, dificultad="Difícil"):
    """
    Genera una sopa de letras (cuadrícula) e intenta colocar todas las palabras
    basado en un nivel de dificultad.
    """
    
    DIRECCIONES_MAP = {
        "Fácil": [
            (0, 1),   # Derecha
            (1, 0)    # Abajo
        ],
        "Medio": [
            (0, 1),   # Derecha
            (1, 0),   # Abajo
            (1, 1),   # Abajo-Derecha
            (1, -1)   # Abajo-Izquierda
        ],
        "Difícil": [
            (1,0), (0,1), (1,1), (-1,1), (-1,0), (0,-1), (-1,-1), (1,-1) # Todas las 8
        ]
    }
    
    direcciones = DIRECCIONES_MAP.get(dificultad, DIRECCIONES_MAP["Difícil"])

    sopa = [[' ' for _ in range(dimension)] for _ in range(dimension)]
    ubicaciones = {}
    palabras_no_colocadas = []
    
    def chequear_fit(palabra, fila, col, dy, dx):
        """Helper interno: Verifica si una palabra cabe en una ubicación."""
        for i in range(len(palabra)):
            fila_actual, col_actual = fila + i * dy, col + i * dx
            if not (0 <= fila_actual < dimension and 0 <= col_actual < dimension):
                return False
            celda = sopa[fila_actual][col_actual]
            if celda != ' ' and celda != palabra[i]:
                return False
        return True

    def colocar_palabra(palabra):
        """Helper interno: Intenta colocar una palabra en la cuadrícula."""
        palabra_upper = palabra.upper()
        posibles_ubicaciones = []
        
        for dy, dx in direcciones:
            for fila_inicio in range(dimension):
                for col_inicio in range(dimension):
                    if chequear_fit(palabra_upper, fila_inicio, col_inicio, dy, dx):
                        posibles_ubicaciones.append((fila_inicio, col_inicio, dy, dx))
        
        if posibles_ubicaciones:
            random.shuffle(posibles_ubicaciones)
            fila_inicio, col_inicio, dy, dx = posibles_ubicaciones[0]
            
            for i in range(len(palabra_upper)):
                fila_actual = fila_inicio + i * dy
                col_actual = col_inicio + i * dx
                sopa[fila_actual][col_actual] = palabra_upper[i]
            
            return (fila_inicio, col_inicio, dx, dy)
        
        return None

    palabras_ordenadas = sorted(palabras, key=len, reverse=True)

    for palabra in palabras_ordenadas:
        ubicacion = colocar_palabra(palabra)
        if ubicacion:
            ubicaciones[palabra.upper()] = ubicacion
        else:
            palabras_no_colocadas.append(palabra)
            
    for i in range(dimension):
        for j in range(dimension):
            if sopa[i][j] == ' ':
                sopa[i][j] = random.choice(LETRAS_ALFABETO)
                
    return sopa, ubicaciones, palabras_no_colocadas

def procesar_excel(excel_file, words_per_puzzle):
    """
    Lee un archivo Excel de dos columnas (Tema, Palabra) y lo divide en 
    listas de palabras.
    """
    try:
        df = pd.read_excel(excel_file, header=None, names=['Tema', 'Palabra'])
        df.dropna(subset=['Tema', 'Palabra'], inplace=True)
        df['Tema'] = df['Tema'].astype(str).str.strip()
        df['Palabra'] = df['Palabra'].astype(str).str.strip()

        word_lists = []
        themes = []
        grouped = df.groupby('Tema')

        for theme_name, group_df in grouped:
            all_words_for_theme = group_df['Palabra'].tolist()
            num_chunks = (len(all_words_for_theme) + words_per_puzzle - 1) // words_per_puzzle
            
            for i in range(0, len(all_words_for_theme), words_per_puzzle):
                chunk = all_words_for_theme[i : i + words_per_puzzle]
                word_lists.append(chunk)
                
                if num_chunks > 1:
                    chunk_num = (i // words_per_puzzle) + 1
                    themes.append(f"{theme_name} ({chunk_num})")
                else:
                    themes.append(theme_name)
                            
        return word_lists, themes

    except Exception as e:
        # Estos mensajes de error se manejan en app_en/es.py
        # st.error(f"Error al procesar el Excel. Verifica el formato.")
        # st.error(f"Asegúrate de que CADA palabra tenga un tema en la columna A.")
        # st.error(f"Detalle: {e}")
        raise e # Relanzar la excepción para que sea capturada por el llamador
        # return [], [] # No es necesario si se relanza

# --- FUNCIONES DE DIBUJO PDF (REPORTLAB) ---

def dibujar_pagina_sopa(c, sopa, palabras, theme, config):
    """
    Dibuja una (1) página completa de sopa de letras en el canvas de ReportLab.
    """
    
    page_width, page_height = config['page_size']
    dimension = config['dimension']
    
    ancho_usable = page_width - 2 * PADDING
    y_cursor = page_height - PADDING

    c.setFont(FONT_NAME_BOLD, TITLE_FONT_SIZE)
    c.drawString(PADDING, y_cursor - TITLE_FONT_SIZE, theme[:60])
    y_cursor -= (TITLE_FONT_SIZE + SPACE_AFTER_TITLE)

    c.setFont(FONT_NAME_REGULAR, WORDS_FONT_SIZE)
    line_height = WORDS_FONT_SIZE * INTERLINEADO_PALABRAS
    
    n_palabras = len(palabras)
    n_cols = NUM_COLUMNAS_PALABRAS
    
    max_words_per_col = (n_palabras + n_cols - 1) // n_cols
    
    col_word_counts = [0] * n_cols
    for i in range(n_palabras):
        col_index = i // max_words_per_col
        if col_index < n_cols:
            col_word_counts[col_index] += 1
            
    font_ascent_guess = WORDS_FONT_SIZE * 0.8

    word_list_height_inner = max_words_per_col * line_height
    word_list_height_total = word_list_height_inner + (2 * BOX_PADDING)

    y_box_bottom = y_cursor - word_list_height_total
    
    c.setLineWidth(1)
    c.setStrokeColor(black)
    c.roundRect(
        PADDING,
        y_box_bottom,
        ancho_usable,
        word_list_height_total,
        BOX_CORNER_RADIUS,
        stroke=1, fill=0
    )

    col_width = (ancho_usable - (2 * BOX_PADDING)) / n_cols
    
    try:
        char_width = pdfmetrics.stringWidth("W", FONT_NAME_REGULAR, WORDS_FONT_SIZE)
        if char_width > 0:
            max_chars = int((col_width * 0.95) / char_width)
        else:
            max_chars = 18
        
        if max_chars < 5: max_chars = 5 

    except Exception as e:
        # Este warning se maneja en app_en/es.py si queremos mostrarlo
        # st.warning(f"Error al calcular ancho de fuente: {e}. Usando fallback de 15 chars.")
        max_chars = 15
    
    for i, word in enumerate(palabras):
        col_index = i // max_words_per_col
        row_index = i % max_words_per_col
        
        words_in_this_col = col_word_counts[col_index]
        this_col_height_inner = words_in_this_col * line_height
        vertical_offset = (word_list_height_inner - this_col_height_inner) / 2
        
        cx = PADDING + BOX_PADDING + (col_index * col_width)
        y_top_of_centered_block = (y_cursor - BOX_PADDING - vertical_offset)
        cy_base = y_top_of_centered_block - font_ascent_guess
        cy = cy_base - (row_index * line_height)

        c.drawString(cx, cy, word.upper()[:max_chars])

    y_cursor -= (word_list_height_total + PADDING)

    grid_available_height = y_cursor - PADDING
    grid_available_width = ancho_usable
    
    cell_size = min(grid_available_width / dimension, grid_available_height / dimension)
    
    grid_width = cell_size * dimension
    grid_height = cell_size * dimension
    
    x_grid = PADDING + (grid_available_width - grid_width) / 2
    y_grid = PADDING + (grid_available_height - grid_height) / 2
    
    c.setLineWidth(1)
    c.setStrokeColor(black)
    c.rect(x_grid, y_grid, grid_width, grid_height, stroke=1, fill=0)
    
    grid_font_size = int(cell_size * PROPORCION_LETRA_CUADRICULA)
    c.setFont(FONT_NAME_REGULAR, grid_font_size)
    
    for i in range(dimension):
        for j in range(dimension):
            letra = sopa[i][j].upper()
            
            x_letra = x_grid + j * cell_size + cell_size / 2
            y_letra = y_grid + grid_height - (i + 0.5) * cell_size
            
            y_letra_ajustada = y_letra - (grid_font_size * 0.3) 
            c.drawCentredString(x_letra, y_letra_ajustada, letra)

def dibujar_paginas_soluciones(c, sopas_list, palabras_list, ubicaciones_list, themes, config, T): # Agregado 'T'
    """
    Dibuja todas las páginas de soluciones, colocando varias (4) por página.
    """
    
    page_width, page_height = config['page_size']
    dimension = config['dimension']

    if not sopas_list:
        return

    c.showPage()
    pagina_solucion_actual = 1
    
    ancho_usable = page_width - 2 * PADDING
    alto_usable = page_height - 2 * PADDING
    
    c.setFont(FONT_NAME_BOLD, TITLE_FONT_SIZE)
    c.drawString(PADDING, page_height - PADDING - TITLE_FONT_SIZE, T['solutions_title']) # Usa T['solutions_title']

    sol_area_width = ancho_usable / 2
    sol_area_height = (alto_usable - (TITLE_FONT_SIZE + PADDING)) / 2
    
    cell_size_sol = min(sol_area_width / dimension, sol_area_height / dimension) * 0.85
    mini_grid_w = cell_size_sol * dimension
    mini_grid_h = cell_size_sol * dimension

    for sol_index, (sopa, palabras, ubicaciones, theme) in enumerate(zip(sopas_list, palabras_list, ubicaciones_list, themes)):
        
        pos_en_pagina = sol_index % SOLUCIONES_POR_PAGINA
        
        if sol_index > 0 and pos_en_pagina == 0:
            c.showPage()
            pagina_solucion_actual += 1
            c.setFont(FONT_NAME_BOLD, TITLE_FONT_SIZE)
            c.drawString(PADDING, page_height - PADDING - TITLE_FONT_SIZE, T['solutions_page_title'].format(page_num=pagina_solucion_actual)) # Usa T['solutions_page_title']

        col_sol = pos_en_pagina % 2
        row_sol = pos_en_pagina // 2
        
        x_offset = PADDING + col_sol * sol_area_width
        left_sol = x_offset + (sol_area_width - mini_grid_w) / 2
        
        y_area_base = PADDING + (1 - row_sol) * sol_area_height
        bot_sol = y_area_base + (sol_area_height - mini_grid_h) / 2
        
        if row_sol == 0:
            bot_sol = (page_height - PADDING - TITLE_FONT_SIZE - PADDING - sol_area_height) + (sol_area_height - mini_grid_h) / 2

        celdas_solucion = set()
        for p in palabras:
            ubicacion = ubicaciones.get(p.upper())
            if ubicacion:
                fila_ini, col_ini, dx, dy = ubicacion
                for k in range(len(p)):
                    celdas_solucion.add((fila_ini + k*dy, col_ini + k*dx))

        c.setFont(FONT_NAME_REGULAR, int(cell_size_sol * 0.7))
        c.setLineWidth(0.5)
        for i in range(dimension):
            for j in range(dimension):
                x = left_sol + j * cell_size_sol
                y = bot_sol + mini_grid_h - (i + 1) * cell_size_sol 
                
                c.rect(x, y, cell_size_sol, cell_size_sol, stroke=1, fill=0)
                
                if (i, j) in celdas_solucion:
                    c.setFillColor(black)
                    c.drawCentredString(x + cell_size_sol / 2, y + (cell_size_sol * 0.2), sopa[i][j].upper())

        c.setFont(FONT_NAME_REGULAR, int(cell_size_sol * 0.6))
        c.setFillColor(black)
        c.drawCentredString(left_sol + mini_grid_w / 2, bot_sol - 14, theme)

def exportar_pdf(sopas_list, palabras_list, ubicaciones_list, themes, dimension, page_size, T): # Agregado 'T'
    """
    Crea el archivo PDF completo en memoria.
    """
    buffer = BytesIO()
    
    try:
        pdfmetrics.registerFont(TTFont(FONT_NAME_REGULAR, TTF_FILE_REGULAR))
        pdfmetrics.registerFont(TTFont(FONT_NAME_BOLD, TTF_FILE_BOLD))
    except Exception as e:
        st.error(T['error_font_load']) # Usa T
        st.error(T['error_font_detail']) # Usa T
        st.error(f"Detail: {e}") # Deja el detalle por si acaso
        return None

    c = canvas.Canvas(buffer, pagesize=page_size)
    
    config = {
        "page_size": page_size,
        "dimension": dimension
    }

    for puzzle_index, (sopa, palabras, ubicaciones, theme) in enumerate(zip(sopas_list, palabras_list, ubicaciones_list, themes)):
        if puzzle_index > 0:
            c.showPage()
        dibujar_pagina_sopa(c, sopa, palabras, theme, config)

    dibujar_paginas_soluciones(c, sopas_list, palabras_list, ubicaciones_list, themes, config, T) # Pasa T

    c.save()
    buffer.seek(0)
    return buffer

# --- FUNCIONES DE GENERACIÓN DE PPTX ---

def dibujar_pagina_sopa_ppt(prs, sopa, palabras, theme, config):
    """
    Dibuja una (1) diapositiva de sopa de letras en la presentación de PPTX.
    """
    page_width, page_height = config['page_size']
    dimension = config['dimension']
    
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    ancho_usable = page_width - (2 * PADDING)
    y_cursor = PADDING

    txBox = slide.shapes.add_textbox(
        Pt(PADDING), Pt(y_cursor), Pt(ancho_usable), Pt(TITLE_FONT_SIZE * 1.2)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = theme[:60]
    p.font.name = FONT_NAME_BOLD
    p.font.size = Pt(TITLE_FONT_SIZE)
    y_cursor += (TITLE_FONT_SIZE + SPACE_AFTER_TITLE)

    n_palabras = len(palabras)
    max_words_per_col = (n_palabras + NUM_COLUMNAS_PALABRAS - 1) // NUM_COLUMNAS_PALABRAS
    
    col_lists = [[] for _ in range(NUM_COLUMNAS_PALABRAS)]
    for i, word in enumerate(palabras):
        col_index = i // max_words_per_col
        if col_index < NUM_COLUMNAS_PALABRAS:
            col_lists[col_index].append(word.upper())
            
    col_width = ancho_usable / NUM_COLUMNAS_PALABRAS
    
    line_height = WORDS_FONT_SIZE * INTERLINEADO_PALABRAS
    word_list_height = (max_words_per_col * line_height) + (2 * BOX_PADDING)

    for col_index in range(NUM_COLUMNAS_PALABRAS):
        text_for_col = "\n".join(col_lists[col_index])
        col_left = PADDING + (col_index * col_width)
        
        txBox_col = slide.shapes.add_textbox(
            Pt(col_left), Pt(y_cursor), Pt(col_width), Pt(word_list_height)
        )
        tf_col = txBox_col.text_frame
        tf_col.text = text_for_col
        
        for p in tf_col.paragraphs:
            p.font.name = FONT_NAME_REGULAR
            p.font.size = Pt(WORDS_FONT_SIZE)
            
    y_cursor += (word_list_height + PADDING)

    grid_available_height = page_height - y_cursor - PADDING
    grid_available_width = ancho_usable
    
    cell_size = min(grid_available_width / dimension, grid_available_height / dimension)
    
    grid_width = cell_size * dimension
    grid_height = cell_size * dimension
    
    x_grid = PADDING + (grid_available_width - grid_width) / 2
    y_grid = y_cursor + (grid_available_height - grid_height) / 2

    shape = slide.shapes.add_table(
        dimension, dimension, Pt(x_grid), Pt(y_grid), Pt(grid_width), Pt(grid_height)
    )
    table = shape.table

    grid_font_size = int(cell_size * PROPORCION_LETRA_CUADRICULA)

    for i in range(dimension):
        table.rows[i].height = Pt(cell_size)
        for j in range(dimension):
            table.columns[j].width = Pt(cell_size)
            
            cell = table.cell(i, j)
            cell.text = sopa[i][j].upper()
            
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.name = FONT_NAME_REGULAR
            p.font.size = Pt(grid_font_size)
            
            cell.margin_top = Pt(0)
            cell.margin_bottom = Pt(0)
            cell.margin_left = Pt(0)
            cell.margin_right = Pt(0)

def dibujar_paginas_soluciones_ppt(prs, sopas_list, palabras_list, ubicaciones_list, themes, config, T): # Agregado 'T'
    """
    Dibuja todas las páginas de soluciones (4 por diapositiva) en PPTX.
    """
    page_width, page_height = config['page_size']
    dimension = config['dimension']
    
    if not sopas_list:
        return
        
    slide = None
    pagina_solucion_actual = 1

    ancho_usable = page_width - 2 * PADDING
    alto_usable = page_height - 2 * PADDING

    sol_area_width = ancho_usable / 2
    sol_area_height = (alto_usable - (TITLE_FONT_SIZE + PADDING)) / 2
    
    cell_size_sol = min(sol_area_width / dimension, sol_area_height / dimension) * 0.85
    mini_grid_w = cell_size_sol * dimension
    mini_grid_h = cell_size_sol * dimension

    for sol_index, (sopa, palabras, ubicaciones, theme) in enumerate(zip(sopas_list, palabras_list, ubicaciones_list, themes)):
        
        pos_en_pagina = sol_index % SOLUCIONES_POR_PAGINA
        
        if pos_en_pagina == 0:
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            title = T['solutions_title'] # Usa T
            if sol_index > 0:
                pagina_solucion_actual += 1
                title = T['solutions_page_title'].format(page_num=pagina_solucion_actual) # Usa T
            
            txBox = slide.shapes.add_textbox(
                Pt(PADDING), Pt(PADDING), Pt(ancho_usable), Pt(TITLE_FONT_SIZE * 1.2)
            )
            txBox.text_frame.paragraphs[0].text = title
            txBox.text_frame.paragraphs[0].font.name = FONT_NAME_BOLD
            txBox.text_frame.paragraphs[0].font.size = Pt(TITLE_FONT_SIZE)

        celdas_solucion = set()
        for p in palabras:
            ubicacion = ubicaciones.get(p.upper())
            if ubicacion:
                fila_ini, col_ini, dx, dy = ubicacion
                for k in range(len(p)):
                    celdas_solucion.add((fila_ini + k*dy, col_ini + k*dx))
                    
        col_sol = pos_en_pagina % 2
        row_sol = pos_en_pagina // 2
        
        x_offset = PADDING + col_sol * sol_area_width
        left_sol = x_offset + (sol_area_width - mini_grid_w) / 2
        
        y_offset = (PADDING + TITLE_FONT_SIZE + PADDING) + row_sol * sol_area_height
        top_sol = y_offset + (sol_area_height - mini_grid_h) / 2

        shape = slide.shapes.add_table(
            dimension, dimension, Pt(left_sol), Pt(top_sol), Pt(mini_grid_w), Pt(mini_grid_h)
        )
        table = shape.table
        sol_font_size = int(cell_size_sol * 0.6)

        for i in range(dimension):
            table.rows[i].height = Pt(cell_size_sol)
            for j in range(dimension):
                table.columns[j].width = Pt(cell_size_sol)
                cell = table.cell(i, j)
                p = cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                p.font.name = FONT_NAME_REGULAR
                p.font.size = Pt(sol_font_size)
                
                cell.margin_top = Pt(0); cell.margin_bottom = Pt(0)
                cell.margin_left = Pt(0); cell.margin_right = Pt(0)
                
                if (i, j) in celdas_solucion:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0, 0, 0)
                    p.text = sopa[i][j].upper()
                    p.font.color.rgb = RGBColor(255, 255, 255)
                else:
                    p.text = ""

        txBox_theme = slide.shapes.add_textbox(
            Pt(left_sol), Pt(top_sol + mini_grid_h + 5), Pt(mini_grid_w), Pt(sol_font_size * 1.5)
        )
        p_theme = txBox_theme.text_frame.paragraphs[0]
        p_theme.text = theme
        p_theme.alignment = PP_ALIGN.CENTER
        p_theme.font.name = FONT_NAME_REGULAR
        p_theme.font.size = Pt(sol_font_size)

def crear_presentacion_ppt(sopas_list, palabras_list, ubicaciones_list, themes, dimension, page_size, T): # Agregado 'T'
    """
    Crea el archivo PPTX completo en memoria.
    """
    buffer = BytesIO()
    prs = Presentation()
    
    prs.slide_width = Pt(page_size[0])
    prs.slide_height = Pt(page_size[1])
    
    config = {
        "page_size": page_size,
        "dimension": dimension
    }

    for sopa, palabras, ubicaciones, theme in zip(sopas_list, palabras_list, ubicaciones_list, themes):
        dibujar_pagina_sopa_ppt(prs, sopa, palabras, theme, config)

    dibujar_paginas_soluciones_ppt(prs, sopas_list, palabras_list, ubicaciones_list, themes, config, T) # Pasa T

    prs.save(buffer)
    buffer.seek(0)
    return buffer