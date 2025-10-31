import streamlit as st
import pandas as pd
import random
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from io import BytesIO
from reportlab.pdfgen import canvas
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.lib.colors import black
from reportlab.lib.units import inch
from reportlab.lib.pagesizes import letter # Usaremos 'letter' como base

# --- CONSTANTES GLOBALES DE DISEÑO ---

# Diccionario de tamaños KDP (ancho, alto en cm)
TAMAÑOS_KDP = {
    '5" x 8" (12.7 x 20.32 cm)': (12.7, 20.32),
    '5.06" x 7.81" (12.85 x 19.84 cm)': (12.85, 19.84),
    '5.25" x 8" (13.34 x 20.32 cm)': (13.34, 20.32),
    '5.5" x 8.5" (13.97 x 21.59 cm)': (13.97, 21.59),
    '6" x 9" (15.24 x 22.86 cm)': (15.24, 22.86),
    '6.14" x 9.21" (15.6 x 23.39 cm)': (15.6, 23.39),
    '6.69" x 9.61" (16.99 x 24.41 cm)': (16.99, 24.41),
    '7" x 10" (17.78 x 25.4 cm)': (17.78, 25.4),
    '8" x 10" (20.32 x 25.4 cm)': (20.32, 25.4),
    '8.25" x 6" (20.96 x 15.24 cm)': (20.96, 15.24),
    '8.25" x 8.25" (20.96 x 20.96 cm)': (20.96, 20.96),
    '8.27" x 11.69" (21.01 x 29.69 cm)': (21.01, 29.69),
    '8.5" x 8.5" (21.59 x 21.59 cm)': (21.59, 21.59),
    '8.5" x 11" (21.59 x 27.94 cm)': (21.59, 27.94)
}

# --- CONSTANTES DE FUENTES ---
FONT_NAME_REGULAR = "RobotoMono-Regular"
FONT_NAME_BOLD = "RobotoMono-Bold"
TTF_FILE_REGULAR = "RobotoMono-Regular.ttf"
TTF_FILE_BOLD = "RobotoMono-Bold.ttf"

# --- CONSTANTES DE DISEÑO Y LAYOUT ---
PADDING = 0.75 * inch              # Margen de seguridad interior (sangrado)
TITLE_FONT_SIZE = 18               # Tamaño fijo para el título de la sopa
WORDS_FONT_SIZE = 11               # Tamaño fijo para la lista de palabras
SPACE_AFTER_TITLE = 0.2 * inch     # Espacio entre título y lista de palabras
INTERLINEADO_PALABRAS = 1.3        # Multiplicador de altura de línea
PROPORCION_LETRA_CUADRICULA = 0.7   # Letra ocupa el 70% de la celda
NUM_COLUMNAS_PALABRAS = 4          # Columnas para la lista de palabras
SOLUCIONES_POR_PAGINA = 4          # Cuadrícula de 2x2 en páginas de solución
LETRAS_ALFABETO = 'ABCDEFGHIJKLMNOPQRSTUVWXYZ'

# Constantes para el recuadro de palabras
BOX_PADDING = 0.15 * inch          # Relleno interno del recuadro
BOX_CORNER_RADIUS = 0.1 * inch     # Radio de las esquinas redondeadas

# --- FUNCIONES DE LÓGICA ---

def cm_to_pt(cm):
    """
    Convierte una medida de centímetros (cm) a puntos (pt) de ReportLab.
    
    Args:
        cm (float): La medida en centímetros.
        
    Returns:
        float: La medida equivalente en puntos.
    """
    return cm * 28.3465

def crear_sopa_letras(palabras, dimension=20, dificultad="Difícil"):
    """
    Genera una sopa de letras (cuadrícula) e intenta colocar todas las palabras
    basado en un nivel de dificultad.
    
    *** Optimización Clave ***:
    Esta función ordena las palabras de más larga a más corta antes de
    intentar colocarlas. Esto aumenta drásticamente la probabilidad de que
    todas las palabras quepan.
    
    Args:
        palabras (list[str]): Lista de palabras a incluir en la sopa.
        dimension (int): Tamaño de la cuadrícula (dimension x dimension).
        dificultad (str): "Fácil", "Medio" o "Difícil". Controla las
                          direcciones permitidas.
        
    Returns:
        tuple: Contiene:
            - sopa (list[list[str]]): La cuadrícula 2D completa con letras.
            - ubicaciones (dict): Diccionario {palabra: (fila, col, dx, dy)}
                                  con la ubicación de las palabras colocadas.
            - palabras_no_colocadas (list[str]): Lista de palabras que
                                                  no se pudieron colocar.
    """
    
    # --- Lógica de Dificultad (¡CORRECTAMENTE INDENTADO!) ---
    # (dy, dx) -> (fila, columna)
    # (1, 0) = Abajo
    # (-1, 0) = Arriba
    # (0, 1) = Derecha
    # (0, -1) = Izquierda
    # (1, 1) = Abajo-Derecha
    # (1, -1) = Abajo-Izquierda
    # (-1, 1) = Arriba-Derecha
    # (-1, -1) = Arriba-Izquierda
    
    DIRECCIONES_MAP = {
        "Fácil": [
            (0, 1),  # Derecha
            (1, 0)   # Abajo
        ],
        "Medio": [
            (0, 1),  # Derecha
            (1, 0),  # Abajo
            (1, 1),  # Abajo-Derecha
            (1, -1)  # Abajo-Izquierda
        ],
        "Difícil": [
            (1,0), (0,1), (1,1), (-1,1), (-1,0), (0,-1), (-1,-1), (1,-1) # Todas las 8
        ]
    }
    
    # Seleccionar las direcciones basadas en la dificultad
    direcciones = DIRECCIONES_MAP.get(dificultad, DIRECCIONES_MAP["Difícil"])
    # --- Fin Lógica de Dificultad ---

    sopa = [[' ' for _ in range(dimension)] for _ in range(dimension)]
    ubicaciones = {}
    palabras_no_colocadas = []
    
    def chequear_fit(palabra, fila, col, dy, dx):
        """Helper interno: Verifica si una palabra cabe en una ubicación."""
        for i in range(len(palabra)):
            fila_actual, col_actual = fila + i * dy, col + i * dx
            # 1. Comprobar límites de la cuadrícula
            if not (0 <= fila_actual < dimension and 0 <= col_actual < dimension):
                return False
            # 2. Comprobar colisiones con otras letras
            celda = sopa[fila_actual][col_actual]
            if celda != ' ' and celda != palabra[i]:
                return False
        return True

    def colocar_palabra(palabra):
        """Helper interno: Intenta colocar una palabra en la cuadrícula."""
        palabra_upper = palabra.upper()
        posibles_ubicaciones = []
        
        # Buscar todas las ubicaciones válidas (usando la lista 'direcciones' seleccionada)
        for dy, dx in direcciones:
            for fila_inicio in range(dimension):
                for col_inicio in range(dimension):
                    if chequear_fit(palabra_upper, fila_inicio, col_inicio, dy, dx):
                        posibles_ubicaciones.append((fila_inicio, col_inicio, dy, dx))
        
        if posibles_ubicaciones:
            # Si hay ubicaciones, elegir una al azar
            random.shuffle(posibles_ubicaciones)
            fila_inicio, col_inicio, dy, dx = posibles_ubicaciones[0]
            
            # Colocar la palabra en la cuadrícula
            for i in range(len(palabra_upper)):
                fila_actual = fila_inicio + i * dy
                col_actual = col_inicio + i * dx
                sopa[fila_actual][col_actual] = palabra_upper[i]
            
            # Devolver la ubicación para las soluciones
            return (fila_inicio, col_inicio, dx, dy)
        
        return None # No se pudo colocar

    # --- ¡OPTIMIZACIÓN! ---
    # Ordenar palabras de más larga a más corta.
    palabras_ordenadas = sorted(palabras, key=len, reverse=True)

    for palabra in palabras_ordenadas:
        ubicacion = colocar_palabra(palabra)
        if ubicacion:
            ubicaciones[palabra.upper()] = ubicacion
        else:
            palabras_no_colocadas.append(palabra)
            
    # Rellenar los espacios vacíos con letras aleatorias
    for i in range(dimension):
        for j in range(dimension):
            if sopa[i][j] == ' ':
                sopa[i][j] = random.choice(LETRAS_ALFABETO)
                
    return sopa, ubicaciones, palabras_no_colocadas

def procesar_excel(excel_file, words_per_puzzle):
    """
    Lee un archivo Excel de dos columnas (Tema, Palabra) y lo divide en 
    listas de palabras.
    
    Formato esperado del Excel:
    - Columna A (índice 0): Tema (ej. "FRUTAS")
    - Columna B (índice 1): Palabra (ej. "Manzana")
    
    El usuario DEBE repetir el tema en la Columna A por cada palabra
    en la Columna B. No debe haber encabezado.
    
    Args:
        excel_file (UploadedFile): El archivo Excel subido desde Streamlit.
        words_per_puzzle (int): El número máximo de palabras para cada sopa.
        
    Returns:
        tuple: Contiene:
            - word_lists (list[list[str]]): Lista de listas de palabras.
            - themes (list[str]): Lista de temas correspondientes.
    """
    try:
        # 1. Leer el Excel. Asumimos que no hay encabezado (header=None).
        #    Nombramos las columnas 'Tema' y 'Palabra' para fácil acceso.
        df = pd.read_excel(excel_file, header=None, names=['Tema', 'Palabra'])
        
        # 2. Limpieza de datos (¡ESTA ES LA LÓGICA CLAVE!)
        #    Eliminar cualquier fila donde falte el Tema O la Palabra.
        #    Esto obliga al usuario a rellenar ambas columnas.
        df.dropna(subset=['Tema', 'Palabra'], inplace=True)

        # 3. Convertir a string y limpiar espacios en blanco (ej. " FRUTAS ")
        df['Tema'] = df['Tema'].astype(str).str.strip()
        df['Palabra'] = df['Palabra'].astype(str).str.strip()

        word_lists = []
        themes = []

        # 4. Agrupar el DataFrame por el 'Tema'
        #    (Ej. todas las filas de "FRUTAS", todas las de "ANIMALES").
        grouped = df.groupby('Tema')

        for theme_name, group_df in grouped:
            # Obtener todas las palabras para este tema como una lista
            all_words_for_theme = group_df['Palabra'].tolist()
            
            # 5. Dividir el grupo grande en "chunks" (trozos)
            #    basado en 'words_per_puzzle'.
            
            # Calcular cuántos chunks saldrán
            num_chunks = (len(all_words_for_theme) + words_per_puzzle - 1) // words_per_puzzle
            
            for i in range(0, len(all_words_for_theme), words_per_puzzle):
                # Cortar la lista de palabras
                chunk = all_words_for_theme[i : i + words_per_puzzle]
                word_lists.append(chunk)
                
                # Asignar el nombre del tema
                if num_chunks > 1:
                    # Si hay más de 1 chunk, añadir un número (ej. "FRUTAS (1)")
                    chunk_num = (i // words_per_puzzle) + 1
                    themes.append(f"{theme_name} ({chunk_num})")
                else:
                    # Si solo hay 1 chunk, usar el nombre del tema tal cual
                    themes.append(theme_name)
                        
        return word_lists, themes

    except Exception as e:
        st.error(f"Error al procesar el Excel. Verifica el formato.")
        st.error(f"Asegúrate de que CADA palabra tenga un tema en la columna A.")
        st.error(f"Detalle: {e}")
        return [], []

# --- FUNCIONES DE DIBUJO PDF (REPORTLAB) ---

def dibujar_pagina_sopa(c, sopa, palabras, theme, config):
    """
    Dibuja una (1) página completa de sopa de letras en el canvas de ReportLab.
    
    Organiza la página en tres secciones (de arriba a abajo):
    1. Título
    2. Recuadro con la lista de palabras
    3. Cuadrícula de la sopa
    
    Args:
        c (canvas.Canvas): El canvas de ReportLab sobre el que dibujar.
        sopa (list[list[str]]): La cuadrícula 2D de letras.
        palabras (list[str]): La lista de palabras a mostrar.
        theme (str): El título de la sopa de letras.
        config (dict): Diccionario con 'page_size' y 'dimension'.
    """
    
    page_width, page_height = config['page_size']
    dimension = config['dimension']
    
    # Área usable (descontando márgenes/padding)
    ancho_usable = page_width - 2 * PADDING
    # 'y_cursor' rastrea nuestra posición vertical, empezando desde arriba
    y_cursor = page_height - PADDING

    # --- 1. Dibujar Título (en negrita) ---
    c.setFont(FONT_NAME_BOLD, TITLE_FONT_SIZE)
    # Dibuja el título en el margen izquierdo, por debajo del padding superior
    c.drawString(PADDING, y_cursor - TITLE_FONT_SIZE, theme[:60])
    # Bajar el cursor
    y_cursor -= (TITLE_FONT_SIZE + SPACE_AFTER_TITLE)

    # --- 2. Dibujar Recuadro y Lista de palabras ---
    c.setFont(FONT_NAME_REGULAR, WORDS_FONT_SIZE)
    line_height = WORDS_FONT_SIZE * INTERLINEADO_PALABRAS
    
    # --- Lógica de Columnas y Centrado Vertical ---
    n_palabras = len(palabras)
    n_cols = NUM_COLUMNAS_PALABRAS
    
    # Calcular cuántas filas tendrá la columna más alta (equivale a math.ceil)
    max_words_per_col = (n_palabras + n_cols - 1) // n_cols
    
    # Contar cuántas palabras reales hay en CADA columna
    col_word_counts = [0] * n_cols
    for i in range(n_palabras):
        col_index = i // max_words_per_col
        if col_index < n_cols:
            col_word_counts[col_index] += 1
            
    # Estimación de la altura de la fuente (para alinear texto verticalmente)
    font_ascent_guess = WORDS_FONT_SIZE * 0.8

    # Calcular altura total del recuadro
    word_list_height_inner = max_words_per_col * line_height
    word_list_height_total = word_list_height_inner + (2 * BOX_PADDING)

    # Coordenadas Y del recuadro
    y_box_bottom = y_cursor - word_list_height_total
    
    # Dibujar el recuadro redondeado
    c.setLineWidth(1)
    c.setStrokeColor(black)
    c.roundRect(
        PADDING,                  # x
        y_box_bottom,             # y
        ancho_usable,             # width
        word_list_height_total,   # height
        BOX_CORNER_RADIUS,        # radius
        stroke=1, fill=0
    )

    # Ancho de cada columna de palabras
    col_width = (ancho_usable - (2 * BOX_PADDING)) / n_cols
    
    # --- ¡NUEVA LÓGICA DE TRUNCAMIENTO! ---
    try:
        # 1. Obtener el ancho de un solo carácter (siendo monospaciada)
        # Usamos 'W' por ser ancha, como medida de seguridad.
        # (Esto funciona porque la fuente se registra en 'exportar_pdf' ANTES de llamar a esta función)
        char_width = pdfmetrics.stringWidth("W", FONT_NAME_REGULAR, WORDS_FONT_SIZE)
        
        # 2. Calcular cuántos caracteres caben en la columna
        # Dejamos un 5% de margen de seguridad para que no se toquen
        if char_width > 0:
             max_chars = int((col_width * 0.95) / char_width)
        else:
            max_chars = 18 # Fallback por si char_width es 0
        
        # Poner un límite mínimo por si acaso
        if max_chars < 5: max_chars = 5 

    except Exception as e:
        # Fallback en caso de que la fuente no esté registrada (no debería pasar)
        st.warning(f"Error al calcular ancho de fuente: {e}. Usando fallback de 15 chars.")
        max_chars = 15 # Usamos un número más seguro que 18
    # --- FIN NUEVA LÓGICA ---
    
    # --- Bucle para dibujar cada palabra ---
    for i, word in enumerate(palabras):
        col_index = i // max_words_per_col
        row_index = i % max_words_per_col
        
        # --- Lógica de Centrado Vertical ---
        # 1. Total de palabras en esta columna específica
        words_in_this_col = col_word_counts[col_index]
        
        # 2. Offset para centrar esta columna (si es más corta que la más alta)
        this_col_height_inner = words_in_this_col * line_height
        vertical_offset = (word_list_height_inner - this_col_height_inner) / 2
        
        # 3. Posición X (columna)
        cx = PADDING + BOX_PADDING + (col_index * col_width)
        
        # 4. Posición Y (fila, ajustada por el centrado)
        y_top_of_centered_block = (y_cursor - BOX_PADDING - vertical_offset)
        cy_base = y_top_of_centered_block - font_ascent_guess
        cy = cy_base - (row_index * line_height)

        # ¡AQUÍ USAMOS EL LÍMITE CALCULADO!
        c.drawString(cx, cy, word.upper()[:max_chars]) # Límite dinámico

    # Bajar el cursor hasta debajo del recuadro de palabras
    y_cursor -= (word_list_height_total + PADDING)

    # --- 3. Dibujar la Cuadrícula (en el espacio restante) ---
    grid_available_height = y_cursor - PADDING
    grid_available_width = ancho_usable
    
    # El tamaño de celda se ajusta al espacio disponible (el menor entre ancho y alto)
    cell_size = min(grid_available_width / dimension, grid_available_height / dimension)
    
    grid_width = cell_size * dimension
    grid_height = cell_size * dimension
    
    # Centrar la cuadrícula en el espacio disponible
    x_grid = PADDING + (grid_available_width - grid_width) / 2
    y_grid = PADDING + (grid_available_height - grid_height) / 2
    
    # Dibujar el borde exterior de la cuadrícula
    c.setLineWidth(1)
    c.setStrokeColor(black)
    c.rect(x_grid, y_grid, grid_width, grid_height, stroke=1, fill=0)
    
    # Tamaño de fuente proporcional al tamaño de la celda
    grid_font_size = int(cell_size * PROPORCION_LETRA_CUADRICULA)
    c.setFont(FONT_NAME_REGULAR, grid_font_size)
    
    # Dibujar cada letra en la cuadrícula
    for i in range(dimension):
        for j in range(dimension):
            letra = sopa[i][j].upper()
            
            # Calcular el centro de la celda
            x_letra = x_grid + j * cell_size + cell_size / 2
            y_letra = y_grid + grid_height - (i + 0.5) * cell_size
            
            # Ajuste vertical fino para centrar la letra (las fuentes no se centran en su 'y' base)
            y_letra_ajustada = y_letra - (grid_font_size * 0.3) 
            c.drawCentredString(x_letra, y_letra_ajustada, letra)

def dibujar_paginas_soluciones(c, sopas_list, palabras_list, ubicaciones_list, themes, config):
    """
    Dibuja todas las páginas de soluciones, colocando varias (4) por página.
    
    Args:
        c (canvas.Canvas): El canvas de ReportLab.
        sopas_list (list): Lista de todas las cuadrículas generadas.
        palabras_list (list): Lista de todas las listas de palabras.
        ubicaciones_list (list): Lista de todos los diccionarios de ubicaciones.
        themes (list): Lista de todos los temas.
        config (dict): Diccionario con 'page_size' y 'dimension'.
    """
    
    page_width, page_height = config['page_size']
    dimension = config['dimension']

    if not sopas_list:
        return # No hay nada que dibujar

    # Iniciar la primera página de soluciones
    c.showPage()
    pagina_solucion_actual = 1
    
    ancho_usable = page_width - 2 * PADDING
    alto_usable = page_height - 2 * PADDING
    
    # Título principal de la sección de soluciones
    c.setFont(FONT_NAME_BOLD, TITLE_FONT_SIZE)
    c.drawString(PADDING, page_height - PADDING - TITLE_FONT_SIZE, "Soluciones")

    # Calcular dimensiones para las mini-sopas (layout 2x2)
    sol_area_width = ancho_usable / 2
    sol_area_height = (alto_usable - (TITLE_FONT_SIZE + PADDING)) / 2
    
    # Calcular tamaño de celda para las mini-cuadrículas
    cell_size_sol = min(sol_area_width / dimension, sol_area_height / dimension) * 0.85
    mini_grid_w = cell_size_sol * dimension
    mini_grid_h = cell_size_sol * dimension

    for sol_index, (sopa, palabras, ubicaciones, theme) in enumerate(zip(sopas_list, palabras_list, ubicaciones_list, themes)):
        
        # Posición dentro de la página actual (0, 1, 2, o 3)
        pos_en_pagina = sol_index % SOLUCIONES_POR_PAGINA
        
        # Si es la primera solución (índice 0) de una NUEVA página...
        if sol_index > 0 and pos_en_pagina == 0:
            c.showPage() # Crear nueva página
            pagina_solucion_actual += 1
            # Añadir título a la nueva página
            c.setFont(FONT_NAME_BOLD, TITLE_FONT_SIZE)
            c.drawString(PADDING, page_height - PADDING - TITLE_FONT_SIZE, f"Soluciones (pág. {pagina_solucion_actual})")

        # Calcular fila (0 o 1) y columna (0 o 1) en la cuadrícula 2x2
        col_sol = pos_en_pagina % 2
        row_sol = pos_en_pagina // 2
        
        # Calcular coordenadas X e Y para esta mini-sopa
        # Centrar la mini-cuadrícula dentro de su cuadrante
        x_offset = PADDING + col_sol * sol_area_width
        left_sol = x_offset + (sol_area_width - mini_grid_w) / 2
        
        # El cálculo de Y es más complejo porque ReportLab empieza desde abajo
        y_area_base = PADDING + (1 - row_sol) * sol_area_height
        bot_sol = y_area_base + (sol_area_height - mini_grid_h) / 2
        
        # Ajuste especial para la fila superior (row_sol == 0) para
        # que quede debajo del título "Soluciones"
        if row_sol == 0:
             bot_sol = (page_height - PADDING - TITLE_FONT_SIZE - PADDING - sol_area_height) + (sol_area_height - mini_grid_h) / 2

        # --- Encontrar todas las celdas que son parte de una solución ---
        celdas_solucion = set()
        for p in palabras:
            ubicacion = ubicaciones.get(p.upper())
            if ubicacion:
                fila_ini, col_ini, dx, dy = ubicacion
                for k in range(len(p)):
                    celdas_solucion.add((fila_ini + k*dy, col_ini + k*dx))

        # --- Dibujar la mini-cuadrícula ---
        c.setFont(FONT_NAME_REGULAR, int(cell_size_sol * 0.7))
        c.setLineWidth(0.5)
        for i in range(dimension):
            for j in range(dimension):
                x = left_sol + j * cell_size_sol
                # 'y' se calcula desde abajo hacia arriba
                y = bot_sol + mini_grid_h - (i + 1) * cell_size_sol 
                
                # Dibujar la celda
                c.rect(x, y, cell_size_sol, cell_size_sol, stroke=1, fill=0)
                
                # Si la celda es una solución, rellenarla
                if (i, j) in celdas_solucion:
                    c.setFillColor(black)
                    c.drawCentredString(x + cell_size_sol / 2, y + (cell_size_sol * 0.2), sopa[i][j].upper())

        # Dibujar el título de la mini-sopa (debajo de ella)
        c.setFont(FONT_NAME_REGULAR, int(cell_size_sol * 0.6))
        c.setFillColor(black)
        c.drawCentredString(left_sol + mini_grid_w / 2, bot_sol - 14, theme)

def exportar_pdf(sopas_list, palabras_list, ubicaciones_list, themes, dimension, page_size):
    """
    Crea el archivo PDF completo en memoria.
    
    Esta función inicializa el canvas, registra las fuentes y luego llama
    a las funciones de dibujo para las páginas de sopas y las de soluciones.
    
    Args:
        sopas_list (list): Lista de todas las cuadrículas.
        palabras_list (list): Lista de todas las listas de palabras.
        ubicaciones_list (list): Lista de todos los diccionarios de ubicaciones.
        themes (list): Lista de todos los temas.
        dimension (int): Tamaño de la cuadrícula.
        page_size (tuple): (ancho, alto) en puntos para el PDF.
        
    Returns:
        BytesIO: Un buffer en memoria que contiene el archivo PDF completo,
                 listo para ser descargado. O None si falla el registro de fuentes.
    """
    buffer = BytesIO()
    
    # --- Registro de Fuentes ---
    # Es crucial registrar las fuentes ANTES de usarlas.
    try:
        pdfmetrics.registerFont(TTFont(FONT_NAME_REGULAR, TTF_FILE_REGULAR))
        pdfmetrics.registerFont(TTFont(FONT_NAME_BOLD, TTF_FILE_BOLD))
    except Exception as e:
        st.error(f"¡ERROR DE FUENTE! No se pudieron cargar los archivos .ttf.")
        st.error(f"Asegúrate de que 'RobotoMono-Regular.ttf' y 'RobotoMono-Bold.ttf' estén en la misma carpeta que el script.")
        st.error(f"Detalle del error: {e}")
        return None # Devuelve None para que Streamlit sepa que falló

    # Crear el canvas (el "lienzo" del PDF)
    c = canvas.Canvas(buffer, pagesize=page_size)
    
    # Crear un diccionario de configuración para pasarlo fácilmente
    config = {
        "page_size": page_size,
        "dimension": dimension
    }

    # --- 1. Dibujar páginas de sopas ---
    for puzzle_index, (sopa, palabras, ubicaciones, theme) in enumerate(zip(sopas_list, palabras_list, ubicaciones_list, themes)):
        if puzzle_index > 0:
            c.showPage() # Crear una nueva página para cada sopa
        dibujar_pagina_sopa(c, sopa, palabras, theme, config)

    # --- 2. Dibujar páginas de soluciones ---
    # Esta función se encarga internamente de crear sus propias páginas
    dibujar_paginas_soluciones(c, sopas_list, palabras_list, ubicaciones_list, themes, config)

    # Guardar y finalizar el PDF
    c.save()
    buffer.seek(0)
    return buffer

# --- FUNCIONES DE GENERACIÓN DE PPTX ---

def dibujar_pagina_sopa_ppt(prs, sopa, palabras, theme, config):
    """
    Dibuja una (1) diapositiva de sopa de letras en la presentación de PPTX.
    
    A diferencia de ReportLab, esto usa formas (cuadros de texto y tablas).
    """
    page_width, page_height = config['page_size']
    dimension = config['dimension']
    
    # 1. Añadir una diapositiva en blanco
    blank_slide_layout = prs.slide_layouts[6] # 'Blank' suele ser el índice 6
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # 2. Definir área usable (en Puntos, ya que page_width/height están en Puntos)
    ancho_usable = page_width - (2 * PADDING)
    y_cursor = PADDING # Empezamos desde el margen superior

    # --- 1. Dibujar Título ---
    # En PPT, añadimos un cuadro de texto
    txBox = slide.shapes.add_textbox(
        Pt(PADDING), Pt(y_cursor), Pt(ancho_usable), Pt(TITLE_FONT_SIZE * 1.2)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = theme[:60]
    p.font.name = FONT_NAME_BOLD
    p.font.size = Pt(TITLE_FONT_SIZE)
    y_cursor += (TITLE_FONT_SIZE + SPACE_AFTER_TITLE)

    # --- 2. Dibujar Lista de palabras ---
    # Usaremos 4 cuadros de texto (uno por columna) para que sea editable
    n_palabras = len(palabras)
    max_words_per_col = (n_palabras + NUM_COLUMNAS_PALABRAS - 1) // NUM_COLUMNAS_PALABRAS
    
    # Dividir palabras en columnas
    col_lists = [[] for _ in range(NUM_COLUMNAS_PALABRAS)]
    for i, word in enumerate(palabras):
        col_index = i // max_words_per_col
        if col_index < NUM_COLUMNAS_PALABRAS:
            col_lists[col_index].append(word.upper())
            
    col_width = ancho_usable / NUM_COLUMNAS_PALABRAS
    
    # Estimar altura de la lista
    line_height = WORDS_FONT_SIZE * INTERLINEADO_PALABRAS
    word_list_height = (max_words_per_col * line_height) + (2 * BOX_PADDING)

    for col_index in range(NUM_COLUMNAS_PALABRAS):
        text_for_col = "\n".join(col_lists[col_index])
        col_left = PADDING + (col_index * col_width)
        
        txBox_col = slide.shapes.add_textbox(
            Pt(col_left), Pt(y_cursor), Pt(col_width), Pt(word_list_height)
        )
        tf_col = txBox_col.text_frame
        tf_col.text = text_for_col
        
        for p in tf_col.paragraphs:
            p.font.name = FONT_NAME_REGULAR
            p.font.size = Pt(WORDS_FONT_SIZE)
            
    y_cursor += (word_list_height + PADDING)

    # --- 3. Dibujar la Cuadrícula (¡como una Tabla!) ---
    grid_available_height = page_height - y_cursor - PADDING
    grid_available_width = ancho_usable
    
    cell_size = min(grid_available_width / dimension, grid_available_height / dimension)
    
    grid_width = cell_size * dimension
    grid_height = cell_size * dimension
    
    # Centrar la cuadrícula
    x_grid = PADDING + (grid_available_width - grid_width) / 2
    y_grid = y_cursor + (grid_available_height - grid_height) / 2

    # Añadir la tabla
    shape = slide.shapes.add_table(
        dimension, dimension, Pt(x_grid), Pt(y_grid), Pt(grid_width), Pt(grid_height)
    )
    table = shape.table

    grid_font_size = int(cell_size * PROPORCION_LETRA_CUADRICULA)

    for i in range(dimension):
        table.rows[i].height = Pt(cell_size) # Definir altura de fila
        for j in range(dimension):
            table.columns[j].width = Pt(cell_size) # Definir ancho de columna
            
            cell = table.cell(i, j)
            cell.text = sopa[i][j].upper()
            
            # Centrar texto en la celda
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.name = FONT_NAME_REGULAR
            p.font.size = Pt(grid_font_size)
            
            # Quitar márgenes internos de la celda
            cell.margin_top = Pt(0)
            cell.margin_bottom = Pt(0)
            cell.margin_left = Pt(0)
            cell.margin_right = Pt(0)

def dibujar_paginas_soluciones_ppt(prs, sopas_list, palabras_list, ubicaciones_list, themes, config):
    """
    Dibuja todas las páginas de soluciones (4 por diapositiva) en PPTX.
    """
    page_width, page_height = config['page_size']
    dimension = config['dimension']
    
    if not sopas_list:
        return
        
    slide = None
    pagina_solucion_actual = 1

    ancho_usable = page_width - 2 * PADDING
    alto_usable = page_height - 2 * PADDING

    # Dimensiones para las mini-sopas (layout 2x2)
    sol_area_width = ancho_usable / 2
    sol_area_height = (alto_usable - (TITLE_FONT_SIZE + PADDING)) / 2
    
    cell_size_sol = min(sol_area_width / dimension, sol_area_height / dimension) * 0.85
    mini_grid_w = cell_size_sol * dimension
    mini_grid_h = cell_size_sol * dimension

    for sol_index, (sopa, palabras, ubicaciones, theme) in enumerate(zip(sopas_list, palabras_list, ubicaciones_list, themes)):
        
        pos_en_pagina = sol_index % SOLUCIONES_POR_PAGINA
        
        # 1. Crear nueva diapositiva si es necesario
        if pos_en_pagina == 0:
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Título de la diapositiva
            title = "Soluciones"
            if sol_index > 0:
                pagina_solucion_actual += 1
                title = f"Soluciones (pág. {pagina_solucion_actual})"
            
            txBox = slide.shapes.add_textbox(
                Pt(PADDING), Pt(PADDING), Pt(ancho_usable), Pt(TITLE_FONT_SIZE * 1.2)
            )
            txBox.text_frame.paragraphs[0].text = title
            txBox.text_frame.paragraphs[0].font.name = FONT_NAME_BOLD
            txBox.text_frame.paragraphs[0].font.size = Pt(TITLE_FONT_SIZE)

        # 2. Calcular celdas de solución
        celdas_solucion = set()
        for p in palabras:
            ubicacion = ubicaciones.get(p.upper())
            if ubicacion:
                fila_ini, col_ini, dx, dy = ubicacion
                for k in range(len(p)):
                    celdas_solucion.add((fila_ini + k*dy, col_ini + k*dx))
                    
        # 3. Calcular posición de la mini-tabla
        col_sol = pos_en_pagina % 2
        row_sol = pos_en_pagina // 2
        
        # Centrar la mini-cuadrícula dentro de su cuadrante
        x_offset = PADDING + col_sol * sol_area_width
        left_sol = x_offset + (sol_area_width - mini_grid_w) / 2
        
        y_offset = (PADDING + TITLE_FONT_SIZE + PADDING) + row_sol * sol_area_height
        top_sol = y_offset + (sol_area_height - mini_grid_h) / 2

        # 4. Dibujar la mini-tabla
        shape = slide.shapes.add_table(
            dimension, dimension, Pt(left_sol), Pt(top_sol), Pt(mini_grid_w), Pt(mini_grid_h)
        )
        table = shape.table
        sol_font_size = int(cell_size_sol * 0.6)

        for i in range(dimension):
            table.rows[i].height = Pt(cell_size_sol)
            for j in range(dimension):
                table.columns[j].width = Pt(cell_size_sol)
                cell = table.cell(i, j)
                p = cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                p.font.name = FONT_NAME_REGULAR
                p.font.size = Pt(sol_font_size)
                
                # Quitar márgenes
                cell.margin_top = Pt(0); cell.margin_bottom = Pt(0)
                cell.margin_left = Pt(0); cell.margin_right = Pt(0)
                
                if (i, j) in celdas_solucion:
                    # Rellenar celda de negro
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0, 0, 0)
                    # Poner texto en blanco
                    p.text = sopa[i][j].upper()
                    p.font.color.rgb = RGBColor(255, 255, 255)
                else:
                    # Opcional: mostrar letras de relleno
                    # p.text = sopa[i][j].upper()
                    # p.font.color.rgb = RGBColor(200, 200, 200) # Gris claro
                    p.text = "" # O dejar en blanco

        # 5. Dibujar el título de la mini-sopa
        txBox_theme = slide.shapes.add_textbox(
            Pt(left_sol), Pt(top_sol + mini_grid_h + 5), Pt(mini_grid_w), Pt(sol_font_size * 1.5)
        )
        p_theme = txBox_theme.text_frame.paragraphs[0]
        p_theme.text = theme
        p_theme.alignment = PP_ALIGN.CENTER
        p_theme.font.name = FONT_NAME_REGULAR
        p_theme.font.size = Pt(sol_font_size)

def crear_presentacion_ppt(sopas_list, palabras_list, ubicaciones_list, themes, dimension, page_size):
    """
    Crea el archivo PPTX completo en memoria.
    """
    buffer = BytesIO()
    prs = Presentation()
    
    # page_size ya está en puntos (pt)
    # Convertir los puntos (float) a EMUs (int) usando el helper Pt()
    prs.slide_width = Pt(page_size[0])
    prs.slide_height = Pt(page_size[1])
    
    config = {
        "page_size": page_size,
        "dimension": dimension
    }

    # 1. Dibujar páginas de sopas
    for sopa, palabras, ubicaciones, theme in zip(sopas_list, palabras_list, ubicaciones_list, themes):
        dibujar_pagina_sopa_ppt(prs, sopa, palabras, theme, config)

    # 2. Dibujar páginas de soluciones
    dibujar_paginas_soluciones_ppt(prs, sopas_list, palabras_list, ubicaciones_list, themes, config)

    prs.save(buffer)
    buffer.seek(0)
    return buffer


# --- INTERFAZ DE STREAMLIT (UI) ---
# Esta sección controla la página web interactiva.

st.title("Generador de Sopas de Letras KDP")
st.write("¡Bienvenido! Sube tu archivo Excel de palabras para empezar.")

# --- Controles de la Barra Lateral ---
# st.sidebar.header("Configuración") # Descomenta si prefieres los controles en la barra lateral

# 1. Carga de archivo
excel_file = st.file_uploader("Selecciona tu archivo Excel", type=["xlsx"])

# 2. Configuración de la Sopa
dimension = st.number_input("Tamaño de cuadrícula (ej. 20x20)", min_value=10, max_value=30, value=20)
words_per_puzzle = st.number_input("Número de palabras por sopa", min_value=5, max_value=40, value=20)

# 3. Nivel de Dificultad (NUEVO)
dificultad = st.selectbox(
    "Nivel de dificultad",
    ("Fácil", "Medio", "Difícil"),
    index=2 # 'Difícil' por defecto
)

# 4. Configuración del PDF
nombre_tamaño = st.selectbox("Tamaño de página KDP", list(TAMAÑOS_KDP.keys()), index=12) # '8.5" x 8.5"' por defecto

formato_salida = st.radio("Elige el formato de salida:", ("PDF", "PPTX"))

# Convertir el tamaño seleccionado a puntos
ancho_cm, alto_cm = TAMAÑOS_KDP[nombre_tamaño]
page_width = cm_to_pt(ancho_cm)
page_height = cm_to_pt(alto_cm)

# --- Lógica de Generación (Botón Principal) ---
if st.button("Generar Libro de Sopas de Letras"):
    if excel_file is not None:
        
        # 1. Procesar el Excel
        word_lists, themes = procesar_excel(excel_file, words_per_puzzle)
        
        if not word_lists:
            st.error("No se encontraron palabras en el archivo Excel. Revisa el formato.")
        else:
            st.info(f"Se van a generar {len(word_lists)} sopas de letras (dificultad: {dificultad}).")
        
            # Listas para almacenar los resultados de cada sopa
            sopas_list = []
            palabras_list_filled = []
            ubicaciones_list = []
            
            # 2. Generar cada Sopa (Lógica)
            with st.spinner("Generando sopas..."):
                for words, theme in zip(word_lists, themes):
                    
                    # PASAR LA DIFICULTAD A LA FUNCIÓN
                    sopa, ubicaciones, no_colocadas = crear_sopa_letras(
                        words, 
                        dimension=dimension, 
                        dificultad=dificultad
                    )
                    
                    # Informar al usuario si algunas palabras no cupieron
                    if no_colocadas:
                        st.warning(f"En la sopa '{theme}', no se pudieron colocar las siguientes palabras: {', '.join(no_colocadas)}")
                    
                    sopas_list.append(sopa)
                    palabras_list_filled.append(words) # Guardamos la lista original de palabras
                    ubicaciones_list.append(ubicaciones)

            # 3. Generar el archivo de salida (PDF o PPTX)
        if formato_salida == "PDF":
            with st.spinner("Creando archivo PDF..."):
                buffer = exportar_pdf(sopas_list, palabras_list_filled, ubicaciones_list, themes, dimension, (page_width, page_height))

            # 4. Ofrecer la descarga
            if buffer: # Si el buffer se creó correctamente
                st.success("¡PDF generado! Descárgalo aquí:")
                st.download_button(
                    label="Descargar PDF",
                    data=buffer,
                    file_name='sopas_de_letras_kdp.pdf',
                    mime='application/pdf'
                )

        elif formato_salida == "PPTX":
            with st.spinner("Creando archivo PPTX (editable)..."):
                buffer = crear_presentacion_ppt(sopas_list, palabras_list_filled, ubicaciones_list, themes, dimension, (page_width, page_height))

            # 4. Ofrecer la descarga
            if buffer: # Si el buffer se creó correctamente
                st.success("¡PPTX generado! Descárgalo aquí:")
                st.download_button(
                    label="Descargar PPTX",
                    data=buffer,
                    file_name='sopas_de_letras_kdp.pptx',
                    mime='application/vnd.openxmlformats-officedocument.presentationml.presentation'
                )
    else:
        st.error("Por favor, sube un archivo Excel.")